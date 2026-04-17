"""
Generate docs/samples/saint-gobain-sample.pptx — a 5-slide demo deck for UC3.

Theme: Saint-Gobain's "SageGlass" dynamic electrochromic glazing — a realistic
product story that works well for a dual-avatar podcast demo (Interviewer +
Expert). Each slide has a title, bullet-friendly content, and speaker notes
that the script-generation prompt can use.

Run: python scripts/generate_sample_pptx.py
Output: docs/samples/saint-gobain-sample.pptx
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

OUT = Path(__file__).resolve().parent.parent / "docs" / "samples" / "saint-gobain-sample.pptx"

# Saint-Gobain blue
SG_BLUE = RGBColor(0x00, 0x33, 0x66)
SG_ACCENT = RGBColor(0x00, 0x9B, 0xDC)
SLATE = RGBColor(0x33, 0x41, 0x55)
CLOUD = RGBColor(0xF1, 0xF5, 0xF9)

SLIDES = [
    {
        "title": "SageGlass — Dynamic Glazing for Net-Zero Buildings",
        "subtitle": "Saint-Gobain Insights · Innovation Podcast",
        "bullets": [
            "Electrochromic glass that tints on demand",
            "Reduces HVAC energy use by up to 20%",
            "A Saint-Gobain product since 2012",
        ],
        "notes": (
            "This is the opening slide. Introduce SageGlass as Saint-Gobain's flagship "
            "dynamic glazing technology. Frame the conversation around decarbonising the "
            "built environment, which is responsible for roughly 40% of global CO2 emissions."
        ),
    },
    {
        "title": "How Electrochromic Tinting Works",
        "bullets": [
            "Five-layer ceramic coating on the inner glass surface",
            "Low-voltage DC pulse migrates lithium ions between layers",
            "Transition from clear (~60% VLT) to deeply tinted (~1% VLT) in minutes",
            "No moving parts — solid-state switching rated for 100,000+ cycles",
        ],
        "notes": (
            "This is the technical depth slide. The expert should explain the five-layer "
            "stack and the ion migration principle. Keep the metaphor accessible — compare "
            "it to a rechargeable battery where the 'charge state' is the tint level."
        ),
    },
    {
        "title": "Why It Matters — Energy & Comfort",
        "bullets": [
            "HVAC savings: 10–20% annual reduction vs. low-e glass",
            "Lighting savings: 60% less artificial lighting near windows",
            "Glare control without blinds or shades",
            "LEED credits: up to 8 under v4.1",
        ],
        "notes": (
            "The interviewer should ask about quantified impact here. Expert shares the key "
            "metrics and grounds them in a relatable example, e.g., a typical 10,000 m² "
            "commercial building saves enough energy to power 40 homes per year."
        ),
    },
    {
        "title": "Deployment Case Study — European HQ Retrofit",
        "bullets": [
            "2,800 m² facade, 1,200 dynamic panes",
            "Installed Q2 2024, fully integrated with BMS",
            "Measured HVAC reduction: 18.4% year-one",
            "Occupant satisfaction survey: +31 NPS",
        ],
        "notes": (
            "Talk about a real-world deployment. Keep it concrete — single building, "
            "specific numbers, measured outcomes. The interviewer should pivot from "
            "'how it works' to 'what it delivers' here."
        ),
    },
    {
        "title": "What's Next — AI-Driven Control",
        "bullets": [
            "ML models predict solar load 15 minutes ahead",
            "Occupant preference learning per zone",
            "Grid-aware tinting during peak pricing",
            "Available in SageGlass SymphonyAI release — 2026",
        ],
        "notes": (
            "Wrap up with the forward-looking story. The expert should describe how AI "
            "turns the glass from a passive product into an active part of the building's "
            "energy strategy. End on a confident, inviting note — this is the customer's "
            "'future' if they partner with Saint-Gobain."
        ),
    },
]


def _add_title_bar(slide, text: str) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(1.1))
    shape.fill.solid()
    shape.fill.fore_color.rgb = SG_BLUE
    shape.line.fill.background()

    tf = shape.text_frame
    tf.margin_left = Inches(0.6)
    tf.margin_top = Inches(0.25)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)


def _add_accent_bar(slide) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(1.1), Inches(13.33), Inches(0.08))
    shape.fill.solid()
    shape.fill.fore_color.rgb = SG_ACCENT
    shape.line.fill.background()


def _add_bullets(slide, bullets: list[str]) -> None:
    box = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.5), Inches(5))
    tf = box.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"•  {b}"
        p.font.size = Pt(22)
        p.font.color.rgb = SLATE
        p.space_after = Pt(14)


def _add_footer(slide, n: int, total: int) -> None:
    box = slide.shapes.add_textbox(Inches(0.4), Inches(7.0), Inches(12.5), Inches(0.4))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = f"Saint-Gobain Insights   ·   Slide {n} of {total}"
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(0x94, 0xA3, 0xB8)


def build() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]

    # Cover slide
    cover = SLIDES[0]
    s = prs.slides.add_slide(blank_layout)
    # Full-bleed blue background
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = SG_BLUE
    bg.line.fill.background()
    # Accent stripe
    stripe = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(4.6), Inches(13.33), Inches(0.06))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = SG_ACCENT
    stripe.line.fill.background()
    # Title
    tbox = s.shapes.add_textbox(Inches(0.8), Inches(2.4), Inches(11.5), Inches(1.8))
    p = tbox.text_frame.paragraphs[0]
    p.text = cover["title"]
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    # Subtitle
    sbox = s.shapes.add_textbox(Inches(0.8), Inches(4.8), Inches(11.5), Inches(0.6))
    p = sbox.text_frame.paragraphs[0]
    p.text = cover["subtitle"]
    p.font.size = Pt(20)
    p.font.color.rgb = SG_ACCENT
    # Notes
    s.notes_slide.notes_text_frame.text = cover["notes"]

    # Content slides
    total = len(SLIDES)
    for i, data in enumerate(SLIDES[1:], start=2):
        s = prs.slides.add_slide(blank_layout)
        bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = CLOUD
        bg.line.fill.background()
        _add_title_bar(s, data["title"])
        _add_accent_bar(s)
        _add_bullets(s, data["bullets"])
        _add_footer(s, i, total)
        s.notes_slide.notes_text_frame.text = data["notes"]

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    return OUT


if __name__ == "__main__":
    out = build()
    print(f"Created: {out} ({out.stat().st_size:,} bytes, {len(SLIDES)} slides)")
