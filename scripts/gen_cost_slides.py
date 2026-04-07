"""Generate RFI Cost Estimation PowerPoint slides.

Cost model: "Batch-First Training Platform"
- Presentations are batch-rendered into videos ONCE (one-time cost per content)
- Learners watch pre-recorded videos unlimited times (zero marginal cost)
- Real-time avatar used only for occasional live Q&A / interactive demos
- Cost scales with CONTENT CREATED, not with number of viewers
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

# ─── Constants ───────────────────────────────────────────────────────
WIDTH = Inches(13.333)
HEIGHT = Inches(7.5)
DARK_BLUE = RGBColor(0x00, 0x33, 0x66)
AZURE_BLUE = RGBColor(0x00, 0x78, 0xD4)
LIGHT_BLUE = RGBColor(0xDE, 0xEC, 0xF9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MED_GRAY = RGBColor(0x66, 0x66, 0x66)
GREEN = RGBColor(0x10, 0x7C, 0x10)
ORANGE = RGBColor(0xCA, 0x50, 0x10)
RED = RGBColor(0xC5, 0x00, 0x00)
ACCENT_TEAL = RGBColor(0x00, 0x80, 0x80)
BG_LIGHT = RGBColor(0xF5, 0xF7, 0xFA)

prs = Presentation()
prs.slide_width = WIDTH
prs.slide_height = HEIGHT


def add_bg(slide, color=BG_LIGHT):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, w, h, fill_color, text="", font_size=12,
              font_color=WHITE, bold=False, align=PP_ALIGN.CENTER,
              shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    s = slide.shapes.add_shape(shape, left, top, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill_color
    s.line.fill.background()
    if text:
        tf = s.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].alignment = align
        run = tf.paragraphs[0].add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.color.rgb = font_color
        run.font.bold = bold
    return s


def add_text(slide, left, top, w, h, text, font_size=14, color=DARK_GRAY,
             bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.bold = bold
    return tb


def add_multi_text(slide, left, top, w, h, lines, default_size=12,
                   default_color=DARK_GRAY):
    tb = slide.shapes.add_textbox(left, top, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, cfg in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = cfg.get("text", "")
        run.font.size = Pt(cfg.get("size", default_size))
        run.font.color.rgb = cfg.get("color", default_color)
        run.font.bold = cfg.get("bold", False)
        p.alignment = cfg.get("align", PP_ALIGN.LEFT)
        p.space_after = Pt(cfg.get("space_after", 4))
    return tb


def add_table(slide, left, top, w, h, rows, cols, data,
              col_widths=None, header_color=AZURE_BLUE):
    tbl_shape = slide.shapes.add_table(rows, cols, left, top, w, h)
    table = tbl_shape.table
    if col_widths:
        for i, cw in enumerate(col_widths):
            table.columns[i].width = cw
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = str(data[r][c]) if r < len(data) and c < len(data[r]) else ""
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(10)
                if r == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = WHITE
                    paragraph.alignment = PP_ALIGN.CENTER
                else:
                    paragraph.font.color.rgb = DARK_GRAY
                    paragraph.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.RIGHT
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
            elif r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xF0, 0xF4, 0xF8)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    return tbl_shape


# ═══════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BLUE)
add_shape(slide, Inches(0), Inches(0), WIDTH, Inches(1.0), AZURE_BLUE,
          shape=MSO_SHAPE.RECTANGLE)
add_text(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.7),
         "MICROSOFT  \u00d7  ACME", 18, WHITE, True)
add_text(slide, Inches(1.5), Inches(2.0), Inches(10), Inches(1.2),
         "RFI \u2014 AI Avatar Solution", 40, WHITE, True)
add_text(slide, Inches(1.5), Inches(3.2), Inches(10), Inches(0.8),
         "Progressive Cost Estimation", 28, RGBColor(0x80, 0xC0, 0xFF))
add_multi_text(slide, Inches(1.5), Inches(4.5), Inches(10), Inches(2.0), [
    {"text": "From Proof of Concept to Enterprise Scale", "size": 16,
     "color": RGBColor(0xCC, 0xDD, 0xEE)},
    {"text": "", "size": 8, "space_after": 12},
    {"text": "Azure AI Services  \u2022  Azure OpenAI  \u2022  VoiceLive Avatar  \u2022  Container Apps",
     "size": 13, "color": RGBColor(0x88, 0xAA, 0xCC)},
    {"text": "", "size": 8, "space_after": 12},
    {"text": "April 2026 \u2014 Pricing estimates based on current Azure published rates",
     "size": 11, "color": RGBColor(0x88, 0x99, 0xAA)},
    {"text": "All prices in USD, West Europe region, subject to enterprise agreement discounts",
     "size": 11, "color": RGBColor(0x88, 0x99, 0xAA)},
])
print("Slide 1: Title")

# ═══════════════════════════════════════════════════════════════════
# SLIDE 2 — COST MODEL PHILOSOPHY
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), WIDTH, Inches(0.9), DARK_BLUE,
          shape=MSO_SHAPE.RECTANGLE)
add_text(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.6),
         'Cost Model: "Batch-First" Training Platform', 24, WHITE, True)

# Key insight box
add_shape(slide, Inches(0.3), Inches(1.2), Inches(12.7), Inches(1.6),
          RGBColor(0xE3, 0xF2, 0xFD), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
add_multi_text(slide, Inches(0.5), Inches(1.3), Inches(12.3), Inches(1.4), [
    {"text": "Key Insight: Cost scales with CONTENT CREATED, not with NUMBER OF VIEWERS",
     "size": 18, "bold": True, "color": DARK_BLUE, "align": PP_ALIGN.CENTER},
    {"text": "", "size": 6, "space_after": 4},
    {"text": "Batch avatar videos are rendered ONCE per presentation, then stored and replayed unlimited times at zero marginal cost.",
     "size": 13, "color": DARK_GRAY, "align": PP_ALIGN.CENTER},
    {"text": "Whether 10 or 10,000 learners watch a video, the Azure cost is the same: $0 per additional view.",
     "size": 13, "color": DARK_GRAY, "align": PP_ALIGN.CENTER},
])

# Two columns: Batch vs Real-time
add_shape(slide, Inches(0.3), Inches(3.1), Inches(6.2), Inches(0.5),
          GREEN, "Batch Video (UC2) \u2014 Primary Delivery", 13, WHITE, True)
add_multi_text(slide, Inches(0.5), Inches(3.7), Inches(5.8), Inches(2.5), [
    {"text": "How it works:", "size": 12, "bold": True, "color": GREEN},
    {"text": "1. Author uploads PowerPoint presentation", "size": 11, "color": DARK_GRAY},
    {"text": "2. System generates avatar video (one-time batch cost)", "size": 11, "color": DARK_GRAY},
    {"text": "3. Video stored in Azure Blob Storage (~$0.018/GB/mo)", "size": 11, "color": DARK_GRAY},
    {"text": "4. Unlimited learners stream/download at near-zero cost", "size": 11, "color": DARK_GRAY},
    {"text": "", "size": 6, "space_after": 4},
    {"text": "Cost per 10-min video: $5.00 (one-time)", "size": 13, "bold": True, "color": GREEN},
    {"text": "Cost per additional viewer: $0.00", "size": 13, "bold": True, "color": GREEN},
])

add_shape(slide, Inches(6.8), Inches(3.1), Inches(6.2), Inches(0.5),
          AZURE_BLUE, "Real-Time Avatar (UC1) \u2014 Interactive Only", 13, WHITE, True)
add_multi_text(slide, Inches(7.0), Inches(3.7), Inches(5.8), Inches(2.5), [
    {"text": "When it's used:", "size": 12, "bold": True, "color": AZURE_BLUE},
    {"text": "1. Live Q&A with avatar over slide content", "size": 11, "color": DARK_GRAY},
    {"text": "2. On-demand interactive exploration of slides", "size": 11, "color": DARK_GRAY},
    {"text": "3. Agent chat for deeper knowledge queries", "size": 11, "color": DARK_GRAY},
    {"text": "4. Typically short sessions (3\u20135 min of avatar)", "size": 11, "color": DARK_GRAY},
    {"text": "", "size": 6, "space_after": 4},
    {"text": "Cost per 5-min Q&A session: $2.50", "size": 13, "bold": True, "color": AZURE_BLUE},
    {"text": "Used by ~10-20% of learners occasionally", "size": 13, "bold": True, "color": AZURE_BLUE},
])

# Bottom note
add_multi_text(slide, Inches(0.3), Inches(6.5), Inches(12.7), Inches(0.8), [
    {"text": "This model makes the solution extremely cost-effective at scale: "
     "50 training videos serving 5,000 learners costs the same as serving 50 learners.",
     "size": 12, "bold": True, "color": ACCENT_TEAL},
])
print("Slide 2: Cost Model Philosophy")

# ═══════════════════════════════════════════════════════════════════
# SLIDE 3 — PROGRESSIVE PHASES OVERVIEW
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), WIDTH, Inches(0.9), DARK_BLUE,
          shape=MSO_SHAPE.RECTANGLE)
add_text(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.6),
         "Progressive Deployment Phases", 24, WHITE, True)

phases = [
    ("Phase 1: PoC / Pilot", "3\u20136 months", [
        "UC1 Silver \u2014 Avatar over PPT",
        "UC2 \u2014 Batch video generation",
        "10 languages (DragonHD voices)",
        "~20 new presentations/month",
        "50 learners, standard avatars",
        "Single Azure region",
    ], GREEN),
    ("Phase 2: Production", "6\u201312 months", [
        "UC1 Basic + Silver + Gold",
        "UC2 + SCORM/LMS output",
        "37 languages (full RFI scope)",
        "~50 new presentations/month",
        "500 learners, SGChat integration",
        "2-region HA deployment",
    ], AZURE_BLUE),
    ("Phase 3: Enterprise", "12+ months", [
        "UC1 + UC2 + UC3 Podcast",
        "Custom avatar (executive)",
        "Full LMS/SCORM ecosystem",
        "~200 new presentations/month",
        "5,000+ learners, global rollout",
        "Multi-region (3+ regions)",
    ], ORANGE),
]

for i, (title, timeline, items, color) in enumerate(phases):
    x = Inches(0.5 + i * 4.2)
    add_shape(slide, x, Inches(1.3), Inches(3.8), Inches(0.7), color,
              title, 16, WHITE, True)
    add_text(slide, x, Inches(2.05), Inches(3.8), Inches(0.35),
             timeline, 12, MED_GRAY)
    for j, item in enumerate(items):
        add_multi_text(slide, x + Inches(0.2), Inches(2.5 + j * 0.38),
                       Inches(3.5), Inches(0.35), [
            {"text": f"\u2192  {item}", "size": 11, "color": DARK_GRAY}
        ])

for i in range(2):
    x = Inches(4.35 + i * 4.2)
    add_text(slide, x, Inches(1.4), Inches(0.5), Inches(0.5),
             "\u25b6", 28, MED_GRAY)

add_multi_text(slide, Inches(0.5), Inches(5.2), Inches(12), Inches(2.0), [
    {"text": "Cost scales with content production, not learner count",
     "size": 14, "bold": True, "color": DARK_BLUE},
    {"text": "The main cost driver is the number of NEW presentations rendered into avatar videos each month.",
     "size": 12, "color": DARK_GRAY},
    {"text": "Once a video is generated, it is stored in Blob Storage and can be viewed by unlimited learners at near-zero cost.",
     "size": 12, "color": MED_GRAY},
    {"text": "Real-time avatar (live Q&A) is used only occasionally and represents a small fraction of the total cost.",
     "size": 12, "color": MED_GRAY},
])
print("Slide 3: Progressive Phases")

# ═══════════════════════════════════════════════════════════════════
# SLIDE 4 — PHASE 1 COST BREAKDOWN (REVISED)
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), WIDTH, Inches(0.9), GREEN,
          shape=MSO_SHAPE.RECTANGLE)
add_text(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.6),
         "Phase 1: PoC / Pilot \u2014 Monthly Cost Estimate", 24, WHITE, True)

add_shape(slide, Inches(0.3), Inches(1.1), Inches(4.8), Inches(2.8),
          RGBColor(0xE8, 0xF5, 0xE9), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
add_multi_text(slide, Inches(0.5), Inches(1.2), Inches(4.4), Inches(2.6), [
    {"text": "Assumptions", "size": 14, "bold": True, "color": GREEN},
    {"text": "\u2022 50 learners (internal Acme pilot)", "size": 11, "color": DARK_GRAY},
    {"text": "\u2022 20 new presentations created/month", "size": 11, "color": DARK_GRAY},
    {"text": "\u2022 Each presentation: ~10 slides, ~10 min video", "size": 11, "color": DARK_GRAY},
    {"text": "\u2022 Batch video: rendered once, watched many times", "size": 11, "color": DARK_GRAY},
    {"text": "\u2022 ~30 live Q&A sessions/month (~5 min each)", "size": 11, "color": DARK_GRAY},
    {"text": "\u2022 10 languages (DragonHD voices)", "size": 11, "color": DARK_GRAY},
    {"text": "\u2022 Standard avatars, single region", "size": 11, "color": DARK_GRAY},
    {"text": "", "size": 6, "space_after": 2},
    {"text": "Learners primarily watch batch videos", "size": 11, "bold": True, "color": GREEN},
])

phase1_data = [
    ["Service", "Usage/Month", "Unit Cost", "Monthly Cost", "% of Total"],
    ["Batch Avatar Synthesis", "200 min (20 videos)", "$0.50/min", "$100", "34.7%"],
    ["VoiceLive Avatar (live Q&A)", "150 min (30 sessions)", "$0.50/min", "$75", "26.0%"],
    ["Azure OpenAI GPT-4.1", "~1M tokens", "$0.002-0.008/1K", "$6", "2.1%"],
    ["OpenAI Embeddings", "~500K tokens", "$0.02/1M", "$0.01", "<0.1%"],
    ["Container Apps (1vCPU, 2GB)", "24/7 after free tier", "$0.000024/s", "$35", "12.1%"],
    ["Cosmos DB Serverless", "~1M RUs + 0.5 GB", "$0.25/1M RU", "$0.40", "<0.1%"],
    ["Blob Storage (Hot LRS)", "~5 GB (videos+slides)", "$0.018/GB", "$0.10", "<0.1%"],
    ["Container Registry (Basic)", "1 registry", "$5/month", "$5", "1.7%"],
    ["Log Analytics", "~1 GB logs", "Free 5GB", "$0", "0%"],
    ["", "", "", "", ""],
    ["TOTAL", "", "", "$222", "100%"],
]

add_table(slide, Inches(5.3), Inches(1.1), Inches(7.7), Inches(4.8),
          len(phase1_data), 5, phase1_data,
          col_widths=[Inches(2.4), Inches(1.3), Inches(1.5), Inches(1.2), Inches(1.3)])

add_multi_text(slide, Inches(0.5), Inches(4.3), Inches(4.4), Inches(3.0), [
    {"text": "Key Insight", "size": 14, "bold": True, "color": GREEN},
    {"text": "", "size": 6, "space_after": 2},
    {"text": "Batch video = one-time cost. A 10-min training video costs $5 to render and can be viewed by all 50 users at $0 extra.",
     "size": 11, "bold": True, "color": DARK_GRAY},
    {"text": "", "size": 6, "space_after": 2},
    {"text": "Live Q&A is optional and short \u2014 only a fraction of users need interactive sessions.",
     "size": 11, "color": MED_GRAY},
    {"text": "", "size": 6, "space_after": 4},
    {"text": "Cost per video produced: ~$5", "size": 13, "bold": True, "color": GREEN},
    {"text": "Cost per learner: ~$4.44/month", "size": 13, "bold": True, "color": GREEN},
    {"text": "", "size": 6, "space_after": 2},
    {"text": "Annual estimate: ~$2,660", "size": 14, "bold": True, "color": DARK_BLUE},
])
print("Slide 4: Phase 1")

# ═══════════════════════════════════════════════════════════════════
# SLIDE 5 — PHASE 2 COST BREAKDOWN (REVISED)
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), WIDTH, Inches(0.9), AZURE_BLUE,
          shape=MSO_SHAPE.RECTANGLE)
add_text(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.6),
         "Phase 2: Production \u2014 Monthly Cost Estimate", 24, WHITE, True)

add_shape(slide, Inches(0.3), Inches(1.1), Inches(4.8), Inches(2.7), LIGHT_BLUE,
          shape=MSO_SHAPE.ROUNDED_RECTANGLE)
add_multi_text(slide, Inches(0.5), Inches(1.2), Inches(4.4), Inches(2.5), [
    {"text": "Assumptions", "size": 14, "bold": True, "color": AZURE_BLUE},
    {"text": "\u2022 500 learners (Acme departments)", "size": 11, "color": DARK_GRAY},
    {"text": "\u2022 50 new presentations/month", "size": 11, "color": DARK_GRAY},
    {"text": "\u2022 Each presentation: ~10 min batch video", "size": 11, "color": DARK_GRAY},
    {"text": "\u2022 ~100 live Q&A sessions/month (~5 min each)", "size": 11, "color": DARK_GRAY},
    {"text": "\u2022 37 languages (full RFI scope)", "size": 11, "color": DARK_GRAY},
    {"text": "\u2022 UC1 Gold: multi-deck AI selection", "size": 11, "color": DARK_GRAY},
    {"text": "\u2022 SGChat + M365 integration", "size": 11, "color": DARK_GRAY},
    {"text": "\u2022 2-region deployment (HA)", "size": 11, "color": DARK_GRAY},
    {"text": "", "size": 6, "space_after": 2},
    {"text": "500 learners watch same batch videos", "size": 11, "bold": True, "color": AZURE_BLUE},
])

phase2_data = [
    ["Service", "Usage/Month", "Unit Cost", "Monthly Cost", "% of Total"],
    ["Batch Avatar Synthesis", "500 min (50 videos)", "$0.50/min", "$250", "19.4%"],
    ["VoiceLive Avatar (live Q&A)", "500 min (100 sessions)", "$0.50/min", "$250", "19.4%"],
    ["Azure OpenAI GPT-4.1", "~5M tokens", "$0.002-0.008/1K", "$25", "1.9%"],
    ["OpenAI Embeddings", "~3M tokens", "$0.02/1M", "$0.06", "<0.1%"],
    ["Azure AI Translator", "~20M chars", "$10/1M", "$200", "15.5%"],
    ["Container Apps (2vCPU, 4GB x2)", "24/7, 2 replicas", "$0.000024/s", "$310", "24.1%"],
    ["Cosmos DB Serverless", "~5M RUs + 5 GB", "$0.25/1M RU", "$3", "0.2%"],
    ["Blob Storage (Hot LRS)", "~30 GB (videos)", "$0.018/GB", "$0.55", "<0.1%"],
    ["Container Registry (Basic)", "1 registry", "$5/month", "$5", "0.4%"],
    ["Log Analytics", "~3 GB logs", "Free 5GB", "$0", "0%"],
    ["AI Search (Basic) \u2014 UC1 Gold", "1 unit", "$73/month", "$73", "5.7%"],
    ["", "", "", "", ""],
    ["TOTAL", "", "", "$1,117", "100%"],
]

add_table(slide, Inches(5.3), Inches(1.1), Inches(7.7), Inches(5.3),
          len(phase2_data), 5, phase2_data,
          col_widths=[Inches(2.4), Inches(1.3), Inches(1.5), Inches(1.2), Inches(1.3)])

add_multi_text(slide, Inches(0.5), Inches(4.2), Inches(4.4), Inches(3.0), [
    {"text": "Key Insight", "size": 14, "bold": True, "color": AZURE_BLUE},
    {"text": "", "size": 6, "space_after": 2},
    {"text": "10x learners (50\u2192500) but only ~5x cost increase because content production scales, not viewership.",
     "size": 12, "bold": True, "color": DARK_GRAY},
    {"text": "", "size": 6, "space_after": 2},
    {"text": "Translation is now a meaningful cost with 37 languages.",
     "size": 11, "color": MED_GRAY},
    {"text": "Compute and infrastructure become a larger share as avatar costs stay modest.",
     "size": 11, "color": MED_GRAY},
    {"text": "", "size": 6, "space_after": 2},
    {"text": "Cost per video produced: ~$5", "size": 13, "bold": True, "color": AZURE_BLUE},
    {"text": "Cost per learner: ~$2.23/month", "size": 13, "bold": True, "color": AZURE_BLUE},
    {"text": "", "size": 6, "space_after": 2},
    {"text": "Annual estimate: ~$13,400", "size": 14, "bold": True, "color": DARK_BLUE},
])
print("Slide 5: Phase 2")

# ═══════════════════════════════════════════════════════════════════
# SLIDE 6 — PHASE 3 COST BREAKDOWN (REVISED)
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), WIDTH, Inches(0.9), ORANGE,
          shape=MSO_SHAPE.RECTANGLE)
add_text(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.6),
         "Phase 3: Enterprise Scale \u2014 Monthly Cost Estimate", 24, WHITE, True)

add_shape(slide, Inches(0.3), Inches(1.1), Inches(4.8), Inches(2.7),
          RGBColor(0xFD, 0xF0, 0xE0), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
add_multi_text(slide, Inches(0.5), Inches(1.2), Inches(4.4), Inches(2.5), [
    {"text": "Assumptions", "size": 14, "bold": True, "color": ORANGE},
    {"text": "\u2022 5,000 learners (global Acme)", "size": 11, "color": DARK_GRAY},
    {"text": "\u2022 200 new presentations/month", "size": 11, "color": DARK_GRAY},
    {"text": "\u2022 Each presentation: ~10 min batch video", "size": 11, "color": DARK_GRAY},
    {"text": "\u2022 ~300 live Q&A sessions/month (~5 min)", "size": 11, "color": DARK_GRAY},
    {"text": "\u2022 UC3: 30 podcasts/month (20 min each)", "size": 11, "color": DARK_GRAY},
    {"text": "\u2022 37 languages, 1 custom avatar model", "size": 11, "color": DARK_GRAY},
    {"text": "\u2022 Full LMS/SCORM integration", "size": 11, "color": DARK_GRAY},
    {"text": "\u2022 Multi-region (3 regions)", "size": 11, "color": DARK_GRAY},
    {"text": "", "size": 6, "space_after": 2},
    {"text": "5,000 learners watch same video library", "size": 11, "bold": True, "color": ORANGE},
])

phase3_data = [
    ["Service", "Usage/Month", "Unit Cost", "Monthly Cost", "% of Total"],
    ["Batch Avatar Synthesis", "2,000 min (200 videos)", "$0.50/min", "$1,000", "14.2%"],
    ["VoiceLive Avatar (live Q&A)", "1,500 min (300 sessions)", "$0.50/min", "$750", "10.7%"],
    ["UC3 Podcast (dual avatar)", "600 min (30 podcasts)", "$0.50/min x2", "$600", "8.5%"],
    ["Custom Avatar Hosting", "1 model, 24/7", "$0.60/hr", "$432", "6.2%"],
    ["Azure OpenAI GPT-4.1", "~30M tokens", "$0.002-0.008/1K", "$150", "2.1%"],
    ["OpenAI Embeddings", "~20M tokens", "$0.02/1M", "$0.40", "<0.1%"],
    ["Azure AI Translator", "~100M chars", "$10/1M", "$1,000", "14.2%"],
    ["Container Apps (4vCPU, 8GB x3)", "24/7, 3 replicas", "$0.000024/s", "$940", "13.4%"],
    ["Cosmos DB Serverless", "~50M RUs + 50 GB", "$0.25/1M RU", "$25", "0.4%"],
    ["Blob Storage (Hot LRS)", "~200 GB (videos)", "$0.018/GB", "$3.60", "<0.1%"],
    ["AI Search (Standard S1)", "1 unit x 3 regions", "$245/unit", "$735", "10.5%"],
    ["Container Registry (Std)", "1 registry", "$20/month", "$20", "0.3%"],
    ["Log Analytics", "~10 GB logs", "$2.76/GB*", "$14", "0.2%"],
    ["Front Door + Key Vault", "Global routing", "~$36/month", "$36", "0.5%"],
    ["", "", "", "", ""],
    ["TOTAL", "", "", "$5,706", "100%"],
]

add_table(slide, Inches(5.3), Inches(1.1), Inches(7.7), Inches(6.0),
          len(phase3_data), 5, phase3_data,
          col_widths=[Inches(2.4), Inches(1.3), Inches(1.5), Inches(1.2), Inches(1.3)])

add_multi_text(slide, Inches(0.5), Inches(4.2), Inches(4.4), Inches(2.8), [
    {"text": "Key Insight", "size": 14, "bold": True, "color": ORANGE},
    {"text": "", "size": 6, "space_after": 2},
    {"text": "100x learners (50\u21925,000) but only ~25x cost because viewers are free.",
     "size": 12, "bold": True, "color": DARK_GRAY},
    {"text": "", "size": 6, "space_after": 2},
    {"text": "Cost is distributed across content production, translation, compute, and search \u2014 no single dominant driver.",
     "size": 11, "color": MED_GRAY},
    {"text": "UC3 Podcast adds modest cost ($600/mo for 30 episodes).",
     "size": 11, "color": MED_GRAY},
    {"text": "", "size": 6, "space_after": 2},
    {"text": "Cost per video produced: ~$5", "size": 13, "bold": True, "color": ORANGE},
    {"text": "Cost per learner: ~$1.14/month", "size": 13, "bold": True, "color": ORANGE},
    {"text": "", "size": 6, "space_after": 2},
    {"text": "Annual estimate: ~$68,500", "size": 14, "bold": True, "color": DARK_BLUE},
])
print("Slide 6: Phase 3")

# ═══════════════════════════════════════════════════════════════════
# SLIDE 7 — PHASE COMPARISON SUMMARY
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), WIDTH, Inches(0.9), DARK_BLUE,
          shape=MSO_SHAPE.RECTANGLE)
add_text(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.6),
         "Progressive Cost Summary \u2014 Phase Comparison", 24, WHITE, True)

summary_data = [
    ["", "Phase 1: PoC", "Phase 2: Production", "Phase 3: Enterprise"],
    ["Learners", "50", "500", "5,000"],
    ["New Videos/Month", "20", "50", "200"],
    ["Use Cases", "UC1 Silver + UC2", "UC1 All + UC2", "UC1 + UC2 + UC3"],
    ["Languages", "10", "37", "37"],
    ["Avatars", "Standard", "Standard", "Standard + Custom"],
    ["Regions", "1", "2", "3"],
    ["", "", "", ""],
    ["Monthly Cost", "$222", "$1,117", "$5,706"],
    ["Annual Cost", "$2,660", "$13,400", "$68,470"],
    ["Per-Learner/Month", "$4.44", "$2.23", "$1.14"],
    ["Per-Video Produced", "$5.00", "$5.00", "$5.00"],
    ["", "", "", ""],
    ["Batch Video %", "45%", "22%", "18%"],
    ["Live Q&A %", "34%", "22%", "11%"],
    ["Infra + AI %", "21%", "56%", "71%"],
]

tbl_shape = add_table(slide, Inches(0.5), Inches(1.2), Inches(12.3), Inches(5.5),
                      len(summary_data), 4, summary_data,
                      col_widths=[Inches(2.5), Inches(3.0), Inches(3.4), Inches(3.4)],
                      header_color=DARK_BLUE)

table = tbl_shape.table
colors_list = [GREEN, AZURE_BLUE, ORANGE]
for r in [8, 9, 10, 11]:
    for c in range(1, 4):
        cell = table.cell(r, c)
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(12)
            p.font.color.rgb = colors_list[c - 1]

print("Slide 7: Comparison")

# ═══════════════════════════════════════════════════════════════════
# SLIDE 8 — COST SCALING / SESSION ANATOMY
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), WIDTH, Inches(0.9), DARK_BLUE,
          shape=MSO_SHAPE.RECTANGLE)
add_text(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.6),
         "Unit Economics: Cost per Content Piece", 24, WHITE, True)

add_shape(slide, Inches(0.3), Inches(1.2), Inches(6.2), Inches(0.5),
          AZURE_BLUE, "Cost to Produce One Training Video", 13, WHITE, True)

session_data = [
    ["Cost Component", "Calculation", "Cost"],
    ["Batch avatar rendering (10 min)", "10 min x $0.50/min", "$5.00"],
    ["GPT-4.1 translation (10 langs)", "~50K chars x $0.002-0.008/1K", "$0.15"],
    ["Embeddings for Q&A index", "~5K tokens x $0.02/1M", "$0.0001"],
    ["Cosmos DB writes", "~10 RUs", "$0.000003"],
    ["Blob Storage (video + slides)", "~200 MB stored/month", "$0.004"],
    ["", "", ""],
    ["TOTAL per video (one-time)", "", "$5.15"],
    ["", "", ""],
    ["Cost per viewer (streaming)", "Video playback from Blob", "~$0.00"],
    ["Cost per Q&A session (live)", "5 min avatar + GPT-4.1", "~$2.55"],
]

add_table(slide, Inches(0.3), Inches(1.8), Inches(6.2), Inches(4.0),
          len(session_data), 3, session_data,
          col_widths=[Inches(2.5), Inches(2.0), Inches(1.7)])

add_shape(slide, Inches(6.8), Inches(1.2), Inches(6.2), Inches(0.5),
          ACCENT_TEAL, "Why This Model is Cost-Effective", 13, WHITE, True)

add_multi_text(slide, Inches(7.0), Inches(1.9), Inches(5.8), Inches(5.0), [
    {"text": "Batch = Invest Once, Serve Forever", "size": 14, "bold": True, "color": ACCENT_TEAL},
    {"text": "A $5 video can train 10,000 employees. Cost per learner approaches $0.",
     "size": 11, "color": DARK_GRAY},
    {"text": "", "size": 6, "space_after": 6},
    {"text": "Compare: Traditional Video Production", "size": 14, "bold": True, "color": ACCENT_TEAL},
    {"text": "Professional training video: $5,000\u2013$15,000 per video with studio, actors, editing.",
     "size": 11, "color": DARK_GRAY},
    {"text": "AI Avatar: $5 per video. That's 1,000x cheaper.",
     "size": 11, "bold": True, "color": DARK_GRAY},
    {"text": "", "size": 6, "space_after": 6},
    {"text": "Update Content Instantly", "size": 14, "bold": True, "color": ACCENT_TEAL},
    {"text": "Update the PPT, re-generate video for $5. No re-shooting, no re-editing.",
     "size": 11, "color": DARK_GRAY},
    {"text": "", "size": 6, "space_after": 6},
    {"text": "Multilingual at Marginal Cost", "size": 14, "bold": True, "color": ACCENT_TEAL},
    {"text": "Generate 37-language versions of a 10-min video: ~$185 total ($5/lang).",
     "size": 11, "color": DARK_GRAY},
    {"text": "Traditional dubbing: $2,000\u2013$5,000 per language.",
     "size": 11, "color": DARK_GRAY},
    {"text": "", "size": 6, "space_after": 6},
    {"text": "Live Q&A is Optional", "size": 14, "bold": True, "color": ACCENT_TEAL},
    {"text": "Most learners watch batch videos. Only ~10-20% use live Q&A, keeping costs low.",
     "size": 11, "color": DARK_GRAY},
])

add_shape(slide, Inches(0.3), Inches(6.0), Inches(12.7), Inches(0.5),
          DARK_BLUE, "Total cost of a 10-min training video in 37 languages: ~$190  |  Traditional equivalent: ~$100,000+",
          12, WHITE, True)
print("Slide 8: Unit Economics")

# ═══════════════════════════════════════════════════════════════════
# SLIDE 9 — OPTIMIZATION SCENARIOS
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), WIDTH, Inches(0.9), DARK_BLUE,
          shape=MSO_SHAPE.RECTANGLE)
add_text(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.6),
         "Cost Optimization Scenarios (Phase 2 \u2014 500 Users)", 24, WHITE, True)

scenarios_data = [
    ["Scenario", "Monthly Cost", "Savings vs Base", "Per-Learner", "Notes"],
    ["Baseline (published rates)", "$1,117", "\u2014", "$2.23", "Phase 2 with 500 learners"],
    ["EA Discount (20% on Avatar)", "$1,017", "-$100 (9%)", "$2.03", "Enterprise agreement pricing"],
    ["Shorter videos (7 min avg)", "$967", "-$150 (13%)", "$1.93", "Optimize presentation length"],
    ["Use Translator instead of GPT", "$917", "-$200 (18%)", "$1.83", "Azure Translator for batch"],
    ["TTS audio-only for Q&A", "$867", "-$250 (22%)", "$1.73", "No avatar for Q&A, voice only"],
    ["Combined optimistic", "$750", "-$367 (33%)", "$1.50", "EA + shorter + audio Q&A"],
]

add_table(slide, Inches(0.3), Inches(1.2), Inches(12.7), Inches(3.5),
          len(scenarios_data), 5, scenarios_data,
          col_widths=[Inches(2.8), Inches(1.5), Inches(2.0), Inches(1.4), Inches(5.0)])

add_shape(slide, Inches(0.3), Inches(5.0), Inches(12.7), Inches(2.2),
          RGBColor(0xE3, 0xF2, 0xFD), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
add_multi_text(slide, Inches(0.5), Inches(5.1), Inches(12.3), Inches(2.0), [
    {"text": "Recommended Cost Strategy", "size": 16, "bold": True, "color": DARK_BLUE},
    {"text": "", "size": 6, "space_after": 4},
    {"text": "1. Negotiate Enterprise Agreement pricing for AI Speech Services (target: 20-30% discount on avatar)",
     "size": 12, "color": DARK_GRAY},
    {"text": '2. Optimize presentation lengths \u2014 shorter, focused modules (7 min) reduce batch rendering costs',
     "size": 12, "color": DARK_GRAY},
    {"text": "3. Use TTS audio-only for Q&A interactions \u2014 reserve avatar rendering for batch videos only",
     "size": 12, "color": DARK_GRAY},
    {"text": "4. Cache all translations on upload \u2014 Cosmos DB cache eliminates repeat translation costs",
     "size": 12, "color": DARK_GRAY},
    {"text": '5. Leverage Azure Cost Management \u2014 set budgets and alerts to track actual vs. estimated usage',
     "size": 12, "color": DARK_GRAY},
])
print("Slide 9: Optimization")

# ═══════════════════════════════════════════════════════════════════
# SLIDE 10 — TCO 3-YEAR VIEW
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), WIDTH, Inches(0.9), DARK_BLUE,
          shape=MSO_SHAPE.RECTANGLE)
add_text(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.6),
         "Total Cost of Ownership (3-Year View)", 24, WHITE, True)

tco_data = [
    ["Cost Category", "Year 1", "Year 2", "Year 3", "3-Year Total"],
    ["", "", "", "", ""],
    ["Phase 1 Azure (6 months)", "$1,330", "\u2014", "\u2014", "$1,330"],
    ["Phase 2 Azure (6mo Y1 + 12mo Y2)", "$6,700", "$13,400", "\u2014", "$20,100"],
    ["Phase 3 Azure (12 months)", "\u2014", "\u2014", "$68,470", "$68,470"],
    ["", "", "", "", ""],
    ["Implementation (one-time)", "", "", "", ""],
    ["Phase 1\u21922 development", "$80,000", "\u2014", "\u2014", "$80,000"],
    ["Phase 2\u21923 development", "\u2014", "$120,000", "\u2014", "$120,000"],
    ["Integration (SGChat, LMS)", "\u2014", "$60,000", "\u2014", "$60,000"],
    ["Custom avatar training", "\u2014", "\u2014", "$30,000", "$30,000"],
    ["", "", "", "", ""],
    ["Support & Operations", "", "", "", ""],
    ["DevOps / SRE (0.5 FTE)", "$50,000", "$50,000", "$50,000", "$150,000"],
    ["", "", "", "", ""],
    ["ANNUAL TOTAL", "$138,030", "$243,400", "$148,470", "$529,900"],
]

tco_tbl = add_table(slide, Inches(0.3), Inches(1.2), Inches(12.7), Inches(5.5),
                    len(tco_data), 5, tco_data,
                    col_widths=[Inches(3.5), Inches(2.0), Inches(2.0), Inches(2.0), Inches(3.2)],
                    header_color=DARK_BLUE)

# Bold the total row
table = tco_tbl.table
last_row = len(tco_data) - 1
for c in range(5):
    cell = table.cell(last_row, c)
    for p in cell.text_frame.paragraphs:
        p.font.bold = True
        p.font.size = Pt(11)
        p.font.color.rgb = DARK_BLUE
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor(0xBB, 0xDE, 0xFB)

add_multi_text(slide, Inches(0.3), Inches(6.8), Inches(12), Inches(0.6), [
    {"text": "* Implementation costs are indicative estimates. Actual costs depend on team composition and project management approach.",
     "size": 9, "color": MED_GRAY},
    {"text": "* Azure costs assume published pay-as-you-go rates. Enterprise Agreement discounts can reduce Azure costs by 20-40%.",
     "size": 9, "color": MED_GRAY},
])
print("Slide 10: TCO")

# ═══════════════════════════════════════════════════════════════════
# SLIDE 11 — AZURE SERVICE UNIT PRICING REFERENCE
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), WIDTH, Inches(0.9), DARK_BLUE,
          shape=MSO_SHAPE.RECTANGLE)
add_text(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.6),
         "Appendix: Azure Service Unit Pricing Reference", 24, WHITE, True)

pricing_data = [
    ["Azure Service", "SKU / Tier", "Unit", "Unit Price", "Notes"],
    ["Batch Avatar Synthesis", "Standard", "per minute", "$0.50",
     "One-time render per video \u2014 main content cost"],
    ["VoiceLive Avatar (real-time)", "Standard", "per minute", "$0.50",
     "Live interactive Q&A only"],
    ["Custom Avatar Hosting", "Custom", "per model/hour", "$0.60",
     "Phase 3 only \u2014 custom executive avatar"],
    ["Azure OpenAI \u2014 GPT-4.1", "S0", "per 1K input tokens", "$0.002",
     "Translation, Q&A, Agent chat"],
    ["Azure OpenAI \u2014 GPT-4.1", "S0", "per 1K output tokens", "$0.008",
     "Responses, generated text"],
    ["OpenAI \u2014 text-embedding-3-small", "S0", "per 1M input tokens", "$0.02",
     "RAG embeddings (negligible cost)"],
    ["Azure AI Translator", "S1", "per 1M characters", "$10.00",
     "Phase 2+ for 37-language expansion"],
    ["Azure Cosmos DB (Serverless)", "Serverless", "per 1M RUs", "$0.25",
     "Metadata, translations cache"],
    ["Azure Cosmos DB (Serverless)", "Serverless", "per GB/month", "$0.25",
     "Storage for slide data"],
    ["Azure Blob Storage", "Standard LRS Hot", "per GB/month", "$0.018",
     "Slide images, PPTX, videos"],
    ["Azure Container Apps", "Consumption", "per vCPU-second", "$0.000024",
     "Free tier: 180K vCPU-sec/mo"],
    ["Azure Container Registry", "Basic", "per month", "$5.00",
     "10 GB included storage"],
    ["Log Analytics", "\u2014", "per GB ingested", "Free 5GB",
     "First 5 GB/month free"],
]

add_table(slide, Inches(0.3), Inches(1.1), Inches(12.7), Inches(5.5),
          len(pricing_data), 5, pricing_data,
          col_widths=[Inches(3.0), Inches(1.5), Inches(2.0), Inches(1.7), Inches(4.5)])

add_text(slide, Inches(0.3), Inches(6.8), Inches(12), Inches(0.4),
         "* All prices USD, West Europe region, April 2026. Subject to change. "
         "Enterprise agreements may include volume discounts.",
         9, MED_GRAY)
print("Slide 11: Pricing Reference")

# ═══════════════════════════════════════════════════════════════════
# SLIDE 11 — ASSUMPTIONS & NEXT STEPS
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), WIDTH, Inches(0.9), DARK_BLUE,
          shape=MSO_SHAPE.RECTANGLE)
add_text(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.6),
         "Assumptions, Notes & Next Steps", 24, WHITE, True)

add_shape(slide, Inches(0.3), Inches(1.2), Inches(6.2), Inches(0.5),
          MED_GRAY, "Key Assumptions", 13, WHITE, True)
add_multi_text(slide, Inches(0.5), Inches(1.8), Inches(5.8), Inches(4.5), [
    {"text": "Usage Assumptions", "size": 12, "bold": True, "color": DARK_BLUE},
    {"text": "\u2022 Average 4 sessions per user per month (Phase 1-2), 3 at scale",
     "size": 10, "color": DARK_GRAY},
    {"text": "\u2022 Average 10 minutes of avatar streaming per session",
     "size": 10, "color": DARK_GRAY},
    {"text": "\u2022 10 slides per presentation, ~500 words per slide",
     "size": 10, "color": DARK_GRAY},
    {"text": "\u2022 2 Q&A questions per session on average",
     "size": 10, "color": DARK_GRAY},
    {"text": "", "size": 6, "space_after": 4},
    {"text": "Pricing Assumptions", "size": 12, "bold": True, "color": DARK_BLUE},
    {"text": "\u2022 Azure published pay-as-you-go rates (April 2026)",
     "size": 10, "color": DARK_GRAY},
    {"text": "\u2022 West Europe region (France Central available)",
     "size": 10, "color": DARK_GRAY},
    {"text": "\u2022 No Enterprise Agreement discounts applied",
     "size": 10, "color": DARK_GRAY},
    {"text": "\u2022 No Azure Reserved Instances or Savings Plans",
     "size": 10, "color": DARK_GRAY},
    {"text": "", "size": 6, "space_after": 4},
    {"text": "Not Included", "size": 12, "bold": True, "color": RED},
    {"text": "\u2022 Network egress / bandwidth costs (typically <1%)",
     "size": 10, "color": DARK_GRAY},
    {"text": "\u2022 Azure Entra ID / Microsoft 365 licensing",
     "size": 10, "color": DARK_GRAY},
    {"text": "\u2022 End-user devices / browsers",
     "size": 10, "color": DARK_GRAY},
    {"text": "\u2022 Project management and change management",
     "size": 10, "color": DARK_GRAY},
])

add_shape(slide, Inches(6.8), Inches(1.2), Inches(6.2), Inches(0.5),
          AZURE_BLUE, "Recommended Next Steps", 13, WHITE, True)
add_multi_text(slide, Inches(7.0), Inches(1.8), Inches(5.8), Inches(4.5), [
    {"text": "1. Validate Usage Assumptions", "size": 12, "bold": True, "color": AZURE_BLUE},
    {"text": "Run a 2-week pilot with 10 users to measure actual session lengths and frequency.",
     "size": 10, "color": DARK_GRAY},
    {"text": "", "size": 6, "space_after": 6},
    {"text": "2. Negotiate EA Pricing", "size": 12, "bold": True, "color": AZURE_BLUE},
    {"text": "Engage Microsoft account team for Acme-specific Azure AI Speech pricing.",
     "size": 10, "color": DARK_GRAY},
    {"text": "", "size": 6, "space_after": 6},
    {"text": "3. Run Azure Pricing Calculator", "size": 12, "bold": True, "color": AZURE_BLUE},
    {"text": "Configure exact workload in Azure Pricing Calculator for region-specific estimates.",
     "size": 10, "color": DARK_GRAY},
    {"text": "", "size": 6, "space_after": 6},
    {"text": "4. Cost Monitoring Strategy", "size": 12, "bold": True, "color": AZURE_BLUE},
    {"text": "Implement Azure Cost Management budgets and alerts from day one.",
     "size": 10, "color": DARK_GRAY},
    {"text": "", "size": 6, "space_after": 6},
    {"text": "5. Phased Go/No-Go Gates", "size": 12, "bold": True, "color": AZURE_BLUE},
    {"text": "Each phase has a cost review gate before scaling to the next tier.",
     "size": 10, "color": DARK_GRAY},
])
print("Slide 12: Assumptions")

# ─── SAVE ─────────────────────────────────────────────────────────
out_path = os.path.join("docs", "RFI-Cost-Estimation.pptx")
prs.save(out_path)
print(f"\nSaved: {out_path}")
print(f"Total slides: {len(prs.slides)}")
