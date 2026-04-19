"""Generate 3 distinct test .pptx files for UC1 multi-deck search testing.

Each deck covers a different topic with topic-specific vocabulary so we can
assert that a topic-specific query returns slides from the correct deck.

Usage:
    python scripts/make_multideck_test.py [output_dir]

Defaults to ./tests/fixtures/uc1/
"""
from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

DECKS = [
    {
        "filename": "climate-action.pptx",
        "language": "en-US",
        "title": "Climate Action & Sustainability",
        "slides": [
            ("Carbon Neutrality Goals",
             "Reaching net-zero emissions by 2050 requires decarbonising the power grid, electrifying transport, and restoring forests.",
             "Emphasise Saint-Gobain's own science-based target and the ACT methodology we follow."),
            ("Renewable Energy Transition",
             "Solar photovoltaics and offshore wind are the fastest-growing sources of renewable electricity worldwide.",
             "Mention the cost curves: solar PV LCOE dropped 89% in a decade."),
            ("Circular Economy for Buildings",
             "Glass, gypsum and insulation materials can be recycled into new products, reducing raw material extraction.",
             "Cite our closed-loop gypsum take-back programme and recycled-content glass wool."),
            ("Carbon Accounting & Scope 3",
             "Scope 3 emissions from purchased goods often dominate a manufacturer's footprint and require supplier engagement.",
             "Reference the GHG Protocol and CDP disclosures."),
            ("Adaptation & Resilience",
             "Rising sea levels, heatwaves and extreme precipitation demand resilient building envelopes and passive cooling.",
             "Point to our cool-roof and high-performance glazing portfolio."),
        ],
    },
    {
        "filename": "ai-ethics.pptx",
        "language": "en-US",
        "title": "Responsible AI & Ethics",
        "slides": [
            ("Principles of Responsible AI",
             "Fairness, transparency, accountability, privacy and safety are the five pillars of responsible AI deployment.",
             "Tie to the EU AI Act risk-based classification."),
            ("Bias in Machine Learning",
             "Training data reflecting historical inequities can cause models to discriminate against protected groups.",
             "Walk through the COMPAS recidivism study and gender-shaded facial analysis results."),
            ("Explainability & Interpretability",
             "SHAP and LIME let practitioners attribute a model's prediction to individual input features.",
             "Contrast model-agnostic vs model-specific explainers."),
            ("Privacy-Preserving ML",
             "Federated learning and differential privacy enable training on sensitive data without centralising it.",
             "Reference Apple's on-device learning and Google's federated keyboard."),
            ("Governance & AI Audits",
             "Model cards, datasheets and third-party audits operationalise AI governance inside an enterprise.",
             "Suggest adopting NIST AI RMF."),
        ],
    },
    {
        "filename": "cloud-security.pptx",
        "language": "en-US",
        "title": "Cloud Security Fundamentals",
        "slides": [
            ("Shared Responsibility Model",
             "In IaaS the provider secures the hypervisor while the customer secures the operating system and data.",
             "Clarify that SaaS shifts most responsibility to the vendor, but identity stays with the customer."),
            ("Identity & Zero Trust",
             "Zero trust assumes breach and verifies every request with strong authentication and least-privilege access.",
             "Reference NIST SP 800-207 and Azure Conditional Access policies."),
            ("Encryption at Rest and in Transit",
             "TLS 1.3 protects data in motion while customer-managed keys in Key Vault protect data at rest.",
             "Recommend disabling TLS 1.0/1.1 and auditing with Defender for Cloud."),
            ("Network Segmentation",
             "Private endpoints, VNet peering and NSGs replace the legacy perimeter with micro-segmented workloads.",
             "Mention Azure Firewall Premium and IDPS signatures."),
            ("Incident Response in the Cloud",
             "Automated playbooks in Sentinel can quarantine compromised identities within minutes of detection.",
             "Highlight MITRE ATT&CK coverage and Kusto hunt queries."),
        ],
    },
    {
        "filename": "medical-devices.pptx",
        "language": "en-US",
        "title": "Medical Device Regulation",
        "slides": [
            ("EU MDR Classification",
             "Medical devices are classified I, IIa, IIb or III depending on invasiveness, duration of contact and active function.",
             "Mention MDR 2017/745 replacing the old directive."),
            ("Clinical Evaluation",
             "Manufacturers must demonstrate clinical safety and performance through literature review and clinical investigations.",
             "Reference MEDDEV 2.7/1 rev 4 guidance."),
            ("Post-Market Surveillance",
             "Vigilance reporting, PSUR and PMCF studies keep real-world evidence flowing back to the manufacturer.",
             "Cite EUDAMED and UDI requirements."),
            ("Quality Management ISO 13485",
             "ISO 13485 requires risk-based processes for design, production and service of medical devices.",
             "Link to ISO 14971 for risk management."),
            ("Cybersecurity for Connected Devices",
             "FDA and MDCG guidance require secure boot, threat modelling and coordinated vulnerability disclosure.",
             "Highlight SBOM requirements."),
        ],
    },
    {
        "filename": "transition-energetique.pptx",
        "language": "fr-FR",
        "title": "La transition énergétique en France",
        "slides": [
            ("Mix énergétique français",
             "Le nucléaire fournit environ 70% de l'électricité française, complété par l'hydraulique et les renouvelables.",
             "Évoquer la fermeture de Fessenheim et la relance du nouveau nucléaire (EPR2)."),
            ("Rénovation énergétique des bâtiments",
             "MaPrimeRénov' et les certificats d'économie d'énergie financent l'isolation thermique des logements anciens.",
             "Citer l'objectif de 700 000 rénovations par an."),
            ("Hydrogène décarboné",
             "L'électrolyse de l'eau alimentée par le nucléaire produit un hydrogène bas-carbone pour l'industrie lourde.",
             "Mentionner la stratégie nationale hydrogène dotée de 9 milliards d'euros."),
            ("Mobilité électrique",
             "Le déploiement des bornes de recharge rapide accélère l'adoption des véhicules électriques.",
             "Référence au règlement AFIR européen."),
            ("Sobriété énergétique",
             "Réduire la consommation avant de produire : c'est le premier levier du plan sobriété énergétique.",
             "Rappeler l'objectif de -10% de consommation d'ici 2024."),
        ],
    },
    {
        "filename": "innovacion-industrial.pptx",
        "language": "es-ES",
        "title": "Innovación industrial y fabricación avanzada",
        "slides": [
            ("Industria 4.0",
             "La digitalización, el IoT industrial y los gemelos digitales transforman la fabricación tradicional.",
             "Mencionar el plan España Digital 2026 y los PERTE industriales."),
            ("Fabricación aditiva",
             "La impresión 3D metálica permite piezas ligeras y geometrías imposibles con mecanizado convencional.",
             "Citar ejemplos aeroespaciales y de herramental."),
            ("Robótica colaborativa",
             "Los cobots trabajan junto a operarios humanos sin jaulas de seguridad, aumentando la productividad.",
             "Referencia a la norma ISO 10218 y TS 15066."),
            ("Economía circular",
             "La remanufactura y el reciclaje avanzado reducen el consumo de materias primas críticas.",
             "Destacar el reglamento europeo sobre diseño ecológico."),
            ("Gemelos digitales",
             "Los modelos virtuales sincronizados con activos físicos optimizan el mantenimiento predictivo y la eficiencia.",
             "Ejemplos en plantas químicas y líneas de vidrio."),
        ],
    },
]


def build_deck(out_path: Path, title: str, slides: list[tuple[str, str, str]]) -> None:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    # Title slide
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = "Saint-Gobain · AI Avatar Presenter · UC1 test fixture"

    # Content slides
    content_layout = prs.slide_layouts[1]
    for title_text, body_text, notes_text in slides:
        s = prs.slides.add_slide(content_layout)
        s.shapes.title.text = title_text
        body_ph = s.placeholders[1]
        tf = body_ph.text_frame
        tf.text = body_text
        for p in tf.paragraphs:
            for r in p.runs:
                r.font.size = Pt(18)
        notes_tf = s.notes_slide.notes_text_frame
        notes_tf.text = notes_text

    prs.save(str(out_path))
    print(f"Wrote {out_path}  ({len(slides) + 1} slides)")


def main() -> None:
    out_dir = Path(sys.argv[1]) if len(sys.argv) > 1 else Path(__file__).resolve().parent.parent / "tests" / "fixtures" / "uc1"
    out_dir.mkdir(parents=True, exist_ok=True)
    for deck in DECKS:
        build_deck(out_dir / deck["filename"], deck["title"], deck["slides"])
    print(f"\nGenerated {len(DECKS)} test decks in {out_dir}")


if __name__ == "__main__":
    main()
