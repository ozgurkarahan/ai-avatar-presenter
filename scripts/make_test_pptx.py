"""Create a tiny 3-slide PPTX for UC2 library testing."""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

OUT = Path(__file__).resolve().parent.parent / "output" / "uc2-test-deck.pptx"
OUT.parent.mkdir(parents=True, exist_ok=True)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDES = [
    ("Frontier Firm in 2026",
     "How AI agents augment every role — from engineers to marketers."),
    ("Three Pillars",
     "1. Your people\n2. Your agents\n3. Your IQ (organizational knowledge)."),
    ("Next Step",
     "Start small: one agent, one business process, measurable outcome."),
]

blank = prs.slide_layouts[6]  # fully blank layout
for i, (title, body) in enumerate(SLIDES, start=1):
    slide = prs.slides.add_slide(blank)
    # dark gradient-ish solid background
    bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)  # 1 = rectangle
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(0x0B, 0x1F, 0x3A)
    bg.line.fill.background()

    # title
    tb = slide.shapes.add_textbox(Inches(0.8), Inches(1.0), Inches(11.5), Inches(2.0))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.runs[0].font.size = Pt(48)
    p.runs[0].font.bold = True
    p.runs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # body
    body_box = slide.shapes.add_textbox(Inches(0.8), Inches(3.2), Inches(11.5), Inches(3.5))
    bf = body_box.text_frame
    bf.word_wrap = True
    for j, line in enumerate(body.split("\n")):
        p = bf.paragraphs[0] if j == 0 else bf.add_paragraph()
        p.text = line
        p.runs[0].font.size = Pt(28)
        p.runs[0].font.color.rgb = RGBColor(0xD9, 0xE2, 0xEC)

    # small footer
    f = slide.shapes.add_textbox(Inches(0.8), Inches(6.8), Inches(11.5), Inches(0.5))
    ff = f.text_frame
    fp = ff.paragraphs[0]
    fp.text = f"Acme · Slide {i} / {len(SLIDES)}"
    fp.runs[0].font.size = Pt(12)
    fp.runs[0].font.color.rgb = RGBColor(0x90, 0xA4, 0xC0)

prs.save(str(OUT))
print(f"wrote {OUT}  ({OUT.stat().st_size/1024:.1f} KB, {len(SLIDES)} slides)")
