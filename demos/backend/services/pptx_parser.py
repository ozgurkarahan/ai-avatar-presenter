"""Parse PPTX files: extract text, notes, and render slides to PNG via LibreOffice."""

from __future__ import annotations

import io
import logging
import subprocess
import tempfile
import uuid
from dataclasses import dataclass, field
from pathlib import Path
from typing import Optional

from pptx import Presentation

log = logging.getLogger(__name__)


@dataclass
class SlideData:
    index: int
    title: str
    body: str
    notes: str
    image_url: Optional[str] = None
    video_url: Optional[str] = None


@dataclass
class PresentationData:
    id: str
    filename: str
    slide_count: int
    slides: list[SlideData] = field(default_factory=list)
    slide_images: list[tuple[int, bytes]] = field(default_factory=list)
    # (slide_index, video_bytes, file_extension e.g. ".mp4")
    slide_videos: list[tuple[int, bytes, str]] = field(default_factory=list)
    pptx_url: Optional[str] = None
    pptx_bytes: Optional[bytes] = None


def parse_pptx(file_bytes: bytes, filename: str, libreoffice_path: str = "soffice") -> PresentationData:
    """Parse a .pptx file, extract slide text/notes, and render slide images."""
    prs = Presentation(io.BytesIO(file_bytes))
    slides: list[SlideData] = []
    slide_videos: list[tuple[int, bytes, str]] = []

    for i, slide in enumerate(prs.slides):
        title = _extract_title(slide)
        body = _extract_body(slide)
        notes = _extract_notes(slide)
        slide_data = SlideData(index=i, title=title, body=body, notes=notes)
        slides.append(slide_data)

        video = _extract_slide_video(slide)
        if video:
            video_bytes, ext = video
            slide_videos.append((i, video_bytes, ext))
            log.info("Extracted embedded video from slide %d (%s, %d bytes)", i + 1, ext, len(video_bytes))

    slide_images = _render_slides(file_bytes, len(slides), libreoffice_path)

    presentation_id = str(uuid.uuid4())
    return PresentationData(
        id=presentation_id,
        filename=filename,
        slide_count=len(slides),
        slides=slides,
        slide_images=slide_images,
        slide_videos=slide_videos,
        pptx_bytes=file_bytes,
    )


def _render_slides(file_bytes: bytes, slide_count: int, libreoffice_path: str) -> list[tuple[int, bytes]]:
    """Render slides to PNG via PowerPoint COM (Windows) or LibreOffice headless."""
    import sys

    with tempfile.TemporaryDirectory() as tmpdir:
        tmpdir_path = Path(tmpdir)
        pptx_path = tmpdir_path / "presentation.pptx"
        pptx_path.write_bytes(file_bytes)

        # Try PowerPoint COM automation on Windows first
        if sys.platform == "win32":
            result = _render_via_powerpoint(pptx_path, tmpdir_path, slide_count)
            if result:
                return result

        # Try LibreOffice headless
        try:
            result = subprocess.run(
                [
                    libreoffice_path, "--headless", "--convert-to", "pdf",
                    "--outdir", str(tmpdir_path), str(pptx_path),
                ],
                capture_output=True,
                timeout=120,
                encoding="utf-8",
                errors="replace",
            )
            log.info("LibreOffice exit code: %d", result.returncode)

            pdf_files = list(tmpdir_path.glob("*.pdf"))
            if pdf_files:
                from pdf2image import convert_from_path
                poppler_path = _find_poppler()
                kwargs = {"dpi": 150, "fmt": "png"}
                if poppler_path:
                    kwargs["poppler_path"] = poppler_path
                images = convert_from_path(str(pdf_files[0]), **kwargs)
                slide_images = []
                for i, img in enumerate(images):
                    buf = io.BytesIO()
                    img.save(buf, format="PNG")
                    slide_images.append((i, buf.getvalue()))
                return slide_images
        except Exception as e:
            log.warning("LibreOffice rendering failed: %s", e)

        # Fallback: placeholder images
        return _fallback_placeholders(slide_count)


def _render_via_powerpoint(pptx_path: Path, output_dir: Path, slide_count: int) -> list[tuple[int, bytes]] | None:
    """Render slides to PNG using PowerPoint COM automation (Windows only)."""
    try:
        import comtypes.client

        powerpoint = comtypes.client.CreateObject("PowerPoint.Application")
        powerpoint.Visible = 1  # Must be visible or minimized for Export to work

        pres = powerpoint.Presentations.Open(str(pptx_path), WithWindow=False)
        try:
            png_dir = output_dir / "png_export"
            png_dir.mkdir(exist_ok=True)
            # Export all slides as PNG (18 = ppSaveAsPNG)
            pres.SaveAs(str(png_dir / "slide"), 18)
        finally:
            pres.Close()

        # PowerPoint SaveAs creates a subdirectory named after the target
        # e.g. SaveAs("png_export/slide", 18) → png_export/slide/Slide1.PNG
        sub_dir = png_dir / "slide"
        slide_images = []
        for i in range(slide_count):
            candidates = [
                sub_dir / f"Slide{i + 1}.PNG",
                sub_dir / f"Slide{i + 1}.png",
                png_dir / f"Slide{i + 1}.PNG",
                png_dir / f"Slide{i + 1}.png",
            ]
            for candidate in candidates:
                if candidate.exists():
                    slide_images.append((i, candidate.read_bytes()))
                    break
            else:
                # Glob inside subdirectory — anchor to avoid Slide10 matching Slide1
                matches = [
                    p for p in sub_dir.glob(f"Slide{i + 1}.*")
                    if p.is_file()
                ] if sub_dir.exists() else []
                if matches:
                    slide_images.append((i, matches[0].read_bytes()))

        if slide_images:
            log.info("Rendered %d slides via PowerPoint COM", len(slide_images))
            return slide_images

        log.warning("PowerPoint COM export produced no images")
        return None
    except Exception as e:
        log.warning("PowerPoint COM rendering failed: %s", e)
        return None


def _find_poppler() -> str | None:
    """Find poppler binaries path on Windows."""
    import shutil
    import sys

    if shutil.which("pdftoppm"):
        return None  # Already on PATH

    if sys.platform == "win32":
        candidates = list(Path(r"C:\tools\poppler").rglob("pdftoppm.exe"))
        if candidates:
            return str(candidates[0].parent)

    return None


def _fallback_placeholders(slide_count: int) -> list[tuple[int, bytes]]:
    """Generate simple placeholder PNGs when no renderer is available."""
    from PIL import Image, ImageDraw, ImageFont

    results = []
    for i in range(slide_count):
        img = Image.new("RGB", (960, 540), color=(240, 240, 240))
        draw = ImageDraw.Draw(img)
        text = f"Slide {i + 1}"
        try:
            font = ImageFont.truetype("LiberationSans-Regular.ttf", 48)
        except OSError:
            try:
                font = ImageFont.truetype("arial.ttf", 48)
            except OSError:
                font = ImageFont.load_default()
        bbox = draw.textbbox((0, 0), text, font=font)
        x = (960 - bbox[2] + bbox[0]) // 2
        y = (540 - bbox[3] + bbox[1]) // 2
        draw.text((x, y), text, fill=(100, 100, 100), font=font)
        buf = io.BytesIO()
        img.save(buf, format="PNG")
        results.append((i, buf.getvalue()))
    return results


# Relationship types used by PowerPoint for embedded/linked video media
_VIDEO_RELTYPES = frozenset({
    "http://schemas.openxmlformats.org/officeDocument/2006/relationships/video",
    "http://schemas.microsoft.com/office/2007/relationships/media",
})

# Content-type prefixes that indicate video
_VIDEO_CT_PREFIXES = ("video/", "application/vnd.ms-asf")

# Extension → MIME type mapping for the file-serving endpoint
_EXT_TO_MIME: dict[str, str] = {
    ".mp4": "video/mp4",
    ".m4v": "video/mp4",
    ".mov": "video/quicktime",
    ".wmv": "video/x-ms-wmv",
    ".avi": "video/x-msvideo",
    ".webm": "video/webm",
}


def _extract_slide_video(slide) -> Optional[tuple[bytes, str]]:
    """Return the first embedded (non-linked) video on a slide as (bytes, ext), or None."""
    for rel in slide.part.rels.values():
        if rel.reltype not in _VIDEO_RELTYPES:
            continue
        # Skip externally linked videos (e.g. YouTube) — only serve embedded files
        if rel.is_external:
            continue
        try:
            part = rel.target_part
            ct: str = part.content_type
            if any(ct.startswith(p) for p in _VIDEO_CT_PREFIXES) or "video" in ct.lower():
                ext = Path(part.partname).suffix.lower() or ".mp4"
                return part.blob, ext
        except Exception as exc:
            log.debug("Skipping video rel %s: %s", rel.reltype, exc)
    return None


def _extract_title(slide) -> str:
    """Extract the title from a slide."""
    if slide.shapes.title and slide.shapes.title.has_text_frame:
        return slide.shapes.title.text_frame.text.strip()
    return ""


def _extract_body(slide) -> str:
    """Extract all body text (non-title) from a slide."""
    texts: list[str] = []
    for shape in slide.shapes:
        if shape == slide.shapes.title:
            continue
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text:
                texts.append(text)
    return "\n".join(texts)


def _extract_notes(slide) -> str:
    """Extract speaker notes from a slide.

    Tries the standard notes_text_frame first, then falls back to
    scanning all shapes in the notes slide for text content.  Some
    PPTX templates store notes in non-placeholder shapes.
    """
    if not slide.has_notes_slide:
        return ""

    notes_slide = slide.notes_slide

    # Primary: standard notes body placeholder
    if notes_slide.notes_text_frame:
        text = notes_slide.notes_text_frame.text.strip()
        if text:
            return text

    # Fallback: collect text from all shapes on the notes slide,
    # skipping the slide-image placeholder (idx 0) which just has
    # a thumbnail label.
    parts: list[str] = []
    for shape in notes_slide.shapes:
        if shape.has_text_frame:
            tf_text = shape.text_frame.text.strip()
            if tf_text and tf_text != slide.shapes.title.text.strip() if slide.shapes.title else True:
                parts.append(tf_text)
    return "\n".join(parts)
