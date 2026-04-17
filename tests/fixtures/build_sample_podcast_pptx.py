"""Generate a sample PPT with speaker notes for UC3 podcast smoke test.

Run: python tests/fixtures/build_sample_podcast_pptx.py
Output: tests/fixtures/saint-gobain-innovation-sample.pptx
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

OUT = Path(__file__).parent / "saint-gobain-innovation-sample.pptx"

SLIDES = [
    {
        "title": "Saint-Gobain: Shaping a Sustainable Future",
        "bullets": [
            "360+ years of materials innovation",
            "Present in 76 countries, 160,000+ employees",
            "Sustainability at the core of every product",
        ],
        "notes": (
            "Welcome to our introduction to Saint-Gobain. "
            "We are a global leader in light and sustainable construction, "
            "operating in 76 countries with more than 160,000 employees. "
            "For over three and a half centuries we have been shaping the habitat and industries of tomorrow, "
            "combining performance, well-being and care for the planet."
        ),
    },
    {
        "title": "Our Four Business Activities",
        "bullets": [
            "High-Performance Solutions for industry",
            "Construction Products (gypsum, insulation, mortars)",
            "Glass for buildings and mobility",
            "Distribution across Europe and Latin America",
        ],
        "notes": (
            "Saint-Gobain is organized around four complementary business activities. "
            "High-Performance Solutions serves industries such as aerospace, healthcare and mobility. "
            "Construction Products delivers gypsum boards, insulation and mortars for energy-efficient buildings. "
            "Our Glass business provides innovative glazing for both buildings and vehicles. "
            "And our Distribution arm brings these solutions to professionals through thousands of branches."
        ),
    },
    {
        "title": "Innovation & R&D at Scale",
        "bullets": [
            "Eight cross-business research centers worldwide",
            "3,700 researchers, 400+ patents filed per year",
            "Partnerships with top universities and start-ups",
        ],
        "notes": (
            "Innovation is the engine of our growth. "
            "We run eight global research centers and collaborate with 100+ universities, "
            "employing 3,700 researchers who file more than 400 patents every year. "
            "Our open innovation approach connects us with start-ups through NOVA External Venturing, "
            "accelerating breakthroughs in low-carbon materials and digital construction."
        ),
    },
    {
        "title": "Grow & Impact: Our 2030 Ambition",
        "bullets": [
            "Net-zero carbon emissions by 2050",
            "-33% CO2 by 2030 vs 2017 baseline",
            "Circular economy across all product lines",
        ],
        "notes": (
            "Our Grow and Impact strategic plan sets the direction for the decade ahead. "
            "We have committed to net-zero carbon emissions by 2050, "
            "with an intermediate target of minus 33 percent CO2 by 2030 compared to our 2017 baseline. "
            "Circular economy principles are being embedded across every product family, "
            "from recycled glass to reusable gypsum boards."
        ),
    },
    {
        "title": "Why Partner with Saint-Gobain?",
        "bullets": [
            "Trusted expertise in materials science",
            "Proven track record on flagship projects",
            "Sustainability metrics embedded in every contract",
        ],
        "notes": (
            "So why should you partner with Saint-Gobain. "
            "Because we combine centuries of materials expertise with a relentless focus on sustainability. "
            "Whether you are building a hospital, a data center or an electric vehicle, "
            "our teams deliver measurable performance and measurable carbon reductions. "
            "Thank you for listening, and we look forward to building the future with you."
        ),
    },
]


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank = prs.slide_layouts[6]  # fully blank layout

    for s in SLIDES:
        slide = prs.slides.add_slide(blank)

        title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12), Inches(1.2))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = s["title"]
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0x0A, 0x2E, 0x5C)

        body_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.8), Inches(5))
        bf = body_box.text_frame
        bf.word_wrap = True
        for i, b in enumerate(s["bullets"]):
            para = bf.paragraphs[0] if i == 0 else bf.add_paragraph()
            para.text = f"\u2022 {b}"
            para.font.size = Pt(22)
            para.font.color.rgb = RGBColor(0x1F, 0x2A, 0x3D)
            para.space_after = Pt(14)

        notes_tf = slide.notes_slide.notes_text_frame
        notes_tf.text = s["notes"]

    prs.save(OUT)
    print(f"Wrote {OUT} ({OUT.stat().st_size} bytes, {len(SLIDES)} slides)")


if __name__ == "__main__":
    build()
