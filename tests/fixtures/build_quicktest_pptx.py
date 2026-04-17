"""Build a *tiny* 2-slide PPT for fast UC3 podcast iteration (~30s total speech)."""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

OUT = Path(__file__).parent / "acme-quicktest.pptx"

SLIDES = [
    {
        "title": "Acme in 30 Seconds",
        "bullets": ["360 years of materials innovation", "76 countries, 160k employees"],
        "notes": "Acme is a global leader in sustainable construction, present in 76 countries with 160,000 employees.",
    },
    {
        "title": "Our 2030 Ambition",
        "bullets": ["Net zero by 2050", "-33% CO2 by 2030"],
        "notes": "Our strategic plan commits us to net-zero carbon emissions by 2050, with a 33 percent reduction in CO2 by 2030.",
    },
]


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    for s in SLIDES:
        slide = prs.slides.add_slide(blank)
        t = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12), Inches(1.2)).text_frame
        t.word_wrap = True
        p = t.paragraphs[0]
        p.text = s["title"]; p.font.size = Pt(36); p.font.bold = True
        p.font.color.rgb = RGBColor(0x0A, 0x2E, 0x5C)

        b = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.8), Inches(5)).text_frame
        b.word_wrap = True
        for i, line in enumerate(s["bullets"]):
            par = b.paragraphs[0] if i == 0 else b.add_paragraph()
            par.text = f"\u2022 {line}"; par.font.size = Pt(22)
            par.font.color.rgb = RGBColor(0x1F, 0x2A, 0x3D); par.space_after = Pt(14)

        slide.notes_slide.notes_text_frame.text = s["notes"]

    prs.save(OUT)
    print(f"Wrote {OUT} ({OUT.stat().st_size} bytes, {len(SLIDES)} slides)")


if __name__ == "__main__":
    build()
