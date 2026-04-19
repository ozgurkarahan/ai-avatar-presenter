"""Generate 9 coherent RFI test decks (3 groups x 3 decks) for UC1/UC2/UC3 e2e tests.

Each deck has 5 slides (title + 3 content + conclusion) with a 2-sentence speaker
note (~300 chars) suitable for avatar narration (~25-30s spoken).

Run:  python tests/fixtures/rfi/_generate.py
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

OUT = Path(__file__).parent

# ---------------------------------------------------------------------------
# Content catalog — 3 groups, 3 decks each, 5 slides each
# Each slide: (title, [bullets], speaker_notes_2_sentences)
# ---------------------------------------------------------------------------

GROUP_A_SAFETY_FR = {
    "securite-bases": {
        "title": "Sécurité sur chantier : les fondamentaux",
        "lang": "fr-FR",
        "tags": ["safety", "construction", "fr"],
        "slides": [
            ("Sécurité sur chantier : les fondamentaux",
             ["Introduction aux règles de base", "Public : nouveaux arrivants", "Durée : 5 minutes"],
             "Bienvenue dans cette formation sur les fondamentaux de la sécurité sur chantier. Nous allons couvrir les trois piliers indispensables : équipement de protection, signalisation et procédures d'urgence."),
            ("Les équipements de protection individuelle",
             ["Casque de chantier obligatoire", "Chaussures de sécurité S3", "Gilet haute visibilité", "Gants adaptés au risque"],
             "L'équipement de protection individuelle est votre première barrière contre les accidents. Casque, chaussures de sécurité et gilet haute visibilité sont obligatoires dès l'entrée sur site."),
            ("Signalisation et balisage",
             ["Panneaux d'interdiction (rouge)", "Panneaux d'obligation (bleu)", "Panneaux d'avertissement (jaune)", "Zones balisées physiquement"],
             "La signalisation sur chantier suit un code couleur universel : rouge pour l'interdiction, bleu pour l'obligation et jaune pour l'avertissement. Respectez toujours les zones balisées physiquement par rubans ou barrières."),
            ("Règles de circulation sur site",
             ["Respecter les voies piétonnes", "Ne jamais passer sous une charge suspendue", "Signaler sa présence aux engins", "Vitesse limitée pour les véhicules"],
             "Sur un chantier, piétons et engins partagent l'espace, ce qui crée des risques majeurs. Respectez les voies piétonnes, ne passez jamais sous une charge suspendue et signalez toujours votre présence aux conducteurs d'engins."),
            ("En résumé : 3 réflexes essentiels",
             ["S'équiper correctement avant d'entrer", "Respecter la signalisation", "Signaler tout risque au chef de chantier"],
             "Retenez trois réflexes simples : équipez-vous correctement avant d'entrer, respectez la signalisation et signalez tout risque au chef de chantier. Votre vigilance protège toute l'équipe."),
        ],
    },
    "risques-chantier": {
        "title": "Identifier et prévenir les risques de chantier",
        "lang": "fr-FR",
        "tags": ["safety", "risk", "fr"],
        "slides": [
            ("Identifier et prévenir les risques de chantier",
             ["Les 4 familles de risques majeurs", "Méthode d'évaluation", "Prérequis : formation sécurité de base"],
             "Après les fondamentaux, nous passons à l'identification des risques spécifiques au chantier. Cette formation vous apprend à reconnaître les situations dangereuses et à agir en prévention."),
            ("Risque de chute",
             ["Chutes de plain-pied (sols glissants)", "Chutes en hauteur (échafaudages)", "Chutes d'objets", "Première cause d'accidents graves"],
             "Le risque de chute représente la première cause d'accidents graves sur chantier. Il se décline en trois types : chute de plain-pied, chute en hauteur et chute d'objets, chacun avec ses mesures préventives."),
            ("Risque lié aux engins et outils",
             ["Engins de levage", "Outils électroportatifs", "Véhicules de chantier", "Zones de manœuvre délimitées"],
             "Les engins et outils sont indispensables mais présentent des risques de collision, coincement ou projection. Respectez les zones de manœuvre délimitées et maintenez toujours une distance de sécurité."),
            ("Risques chimiques et poussières",
             ["Solvants et produits volatils", "Poussières de silice et d'amiante", "Masque FFP2 ou FFP3 selon produit", "Fiche de sécurité obligatoire"],
             "Les risques chimiques et poussières sont souvent invisibles mais peuvent causer des pathologies graves à long terme. Consultez toujours la fiche de sécurité du produit et portez le masque adapté."),
            ("Méthode : observer, évaluer, agir",
             ["Observer l'environnement avant d'intervenir", "Évaluer les risques associés", "Agir : supprimer, réduire ou protéger"],
             "Adoptez systématiquement la méthode observer-évaluer-agir avant toute intervention. L'objectif est de supprimer le risque à la source, ou à défaut de le réduire ou de s'en protéger."),
        ],
    },
    "intervention-urgence": {
        "title": "Intervention en cas d'urgence",
        "lang": "fr-FR",
        "tags": ["safety", "emergency", "fr"],
        "slides": [
            ("Intervention en cas d'urgence sur chantier",
             ["Les 3 types d'urgence majeurs", "Chaîne des secours", "Prérequis : formation risques chantier"],
             "Cette dernière formation du parcours sécurité couvre l'intervention en situation d'urgence. Vous apprendrez à réagir face à un accident, un incendie ou une évacuation."),
            ("Protéger, alerter, secourir",
             ["Protéger : sécuriser la zone", "Alerter : 18 ou chef de chantier", "Secourir : selon formation SST", "Ne jamais improviser"],
             "Face à un accident, appliquez la règle protéger-alerter-secourir. Protégez d'abord la zone pour éviter un sur-accident, alertez les secours, puis secourez uniquement si vous avez la formation SST."),
            ("Incendie : actions immédiates",
             ["Donner l'alerte immédiatement", "Évacuer calmement par les issues", "Utiliser l'extincteur adapté si formé", "Point de rassemblement signalé"],
             "En cas d'incendie, l'alerte et l'évacuation priment sur toute tentative d'extinction. N'utilisez l'extincteur que si vous êtes formé et que le feu est naissant, puis rejoignez le point de rassemblement signalé."),
            ("Évacuation d'urgence",
             ["Suivre les flèches vertes", "Ne pas revenir en arrière", "Laisser les affaires personnelles", "Rejoindre le point de rassemblement"],
             "Lors d'une évacuation d'urgence, suivez les flèches vertes et ne revenez jamais en arrière chercher vos affaires. Le point de rassemblement permet de vérifier que toute l'équipe est en sécurité."),
            ("En résumé : être prêt avant l'urgence",
             ["Connaître les issues de secours", "Mémoriser les numéros d'urgence", "Participer aux exercices d'évacuation"],
             "La préparation fait toute la différence le jour où l'urgence survient. Connaissez les issues, mémorisez les numéros d'urgence et participez activement aux exercices d'évacuation."),
        ],
    },
}

GROUP_B_SUSTAINABILITY_EN = {
    "decarbonization-intro": {
        "title": "Decarbonization: Introduction for industry",
        "lang": "en-US",
        "tags": ["sustainability", "climate", "en"],
        "slides": [
            ("Decarbonization: An introduction for industry",
             ["Why decarbonization matters", "Scope 1, 2 and 3 emissions", "Target audience: all employees"],
             "Welcome to this introduction on decarbonization for industrial companies. We will cover why it matters, how emissions are categorized, and what each of us can do to contribute."),
            ("Climate science in three facts",
             ["Global temperature +1.2°C since 1900", "CO2 concentration now 420 ppm", "Industry = 30% of global emissions", "Paris Agreement target: limit to +1.5°C"],
             "Three facts frame the urgency: global temperatures have already risen by 1.2 degrees, CO2 concentration is at a three-million-year high, and industry accounts for thirty percent of emissions."),
            ("Scope 1, 2 and 3 emissions explained",
             ["Scope 1: direct emissions (factories, vehicles)", "Scope 2: indirect from purchased electricity", "Scope 3: value chain (suppliers, products in use)", "Scope 3 is often the largest"],
             "Corporate emissions are reported in three scopes. Scope 1 covers what you burn directly, scope 2 the electricity you buy, and scope 3 everything else in your value chain, which is usually the largest share."),
            ("Levers available to companies",
             ["Energy efficiency in plants", "Switch to renewable electricity", "Low-carbon raw materials", "Product design for circularity"],
             "Companies have four main decarbonization levers: improving energy efficiency, switching to renewable electricity, sourcing low-carbon raw materials, and redesigning products for circularity."),
            ("What comes next in this path",
             ["Deep dive: low-carbon materials", "Deep dive: circular economy", "Then: action plan for your team"],
             "This introduction is the first step of a three-deck path. The next decks go deeper into low-carbon materials and circular economy, with concrete actions your team can take."),
        ],
    },
    "low-carbon-materials": {
        "title": "Low-carbon materials in construction",
        "lang": "en-US",
        "tags": ["sustainability", "materials", "en"],
        "slides": [
            ("Low-carbon materials in construction",
             ["Why materials matter: 40% of built-environment CO2", "Main categories", "Prerequisite: decarbonization intro"],
             "Materials represent around forty percent of built-environment emissions, making them a critical lever. This deck reviews the main low-carbon materials available today and their trade-offs."),
            ("Cement and concrete innovations",
             ["Clinker substitution (slag, fly ash)", "Limestone calcined clay (LC3)", "CO2-cured concrete", "Up to 60% emission reduction"],
             "Cement is the largest single source of industrial CO2. Innovations like clinker substitution and LC3 technology can cut emissions by up to sixty percent without changing construction practices."),
            ("Bio-based and recycled insulation",
             ["Wood fiber, hemp, cellulose", "Recycled cotton and glass wool", "Lower embodied carbon", "Often equal or better thermal performance"],
             "Insulation is another high-impact area where bio-based materials like wood fiber or hemp can replace oil-derived foams. They often match or exceed thermal performance while storing carbon in the building."),
            ("Recycled steel and aluminum",
             ["Recycled steel: 75% less energy than virgin", "Recycled aluminum: 95% less", "Mature industrial supply chains", "Requires end-of-life collection"],
             "Recycled metals offer dramatic savings: recycled steel uses seventy-five percent less energy than virgin, and recycled aluminum saves ninety-five percent. The challenge is organizing end-of-life collection at scale."),
            ("How to specify low-carbon materials",
             ["Use Environmental Product Declarations (EPDs)", "Compare per functional unit", "Check local availability", "Pilot on one project first"],
             "To specify low-carbon materials, start from Environmental Product Declarations and compare on a functional-unit basis. Pilot on a single project before rolling out, and always check local supply to avoid transport emissions."),
        ],
    },
    "circular-economy": {
        "title": "Circular economy in practice",
        "lang": "en-US",
        "tags": ["sustainability", "circular", "en"],
        "slides": [
            ("Circular economy in practice",
             ["From linear to circular models", "Four circular strategies", "Prerequisite: low-carbon materials"],
             "The final deck of this sustainability path shifts from reducing emissions to rethinking the whole value chain. Circular economy turns waste into resource and extends product life."),
            ("From linear to circular thinking",
             ["Linear: take-make-waste", "Circular: reduce-reuse-recycle-regenerate", "Keep materials at highest value", "Design for disassembly"],
             "The linear take-make-waste model reaches its limits as resources get scarcer. Circular economy keeps materials at their highest value through reduce, reuse, recycle and regenerate strategies."),
            ("Reuse before recycle",
             ["Component reuse: windows, beams", "Refurbishment programs", "Second-life markets", "Avoids reprocessing energy"],
             "Reuse almost always beats recycling because it avoids the energy of reprocessing. Component reuse programs for windows, beams or equipment can extend lifespan and unlock second-life markets."),
            ("Industrial symbiosis",
             ["One company's waste = another's feedstock", "Example: slag from steel to cement", "Regional clusters work best", "Shared infrastructure key"],
             "Industrial symbiosis turns one company's waste into another's feedstock, like steel slag feeding cement plants. Regional clusters with shared infrastructure make the economics work."),
            ("Getting started: 3 quick wins",
             ["Audit waste streams on one site", "Find a local reuse partner", "Set a 10% recycled-content target"],
             "To start, audit waste streams at one site to see the value leaking out. Find a local reuse partner and set an achievable recycled-content target; small wins build momentum for bigger changes."),
        ],
    },
}

GROUP_C_AI_ES = {
    "ia-industrial-intro": {
        "title": "Introducción a la IA en la industria",
        "lang": "es-ES",
        "tags": ["ai", "industry", "es"],
        "slides": [
            ("Introducción a la IA en la industria",
             ["Qué es la IA aplicada a la industria", "Casos de uso principales", "Destinatarios: todo el personal"],
             "Bienvenidos a esta introducción a la inteligencia artificial en la industria. Veremos qué es la IA aplicada, los casos de uso principales y cómo se integra en los procesos productivos."),
            ("Tres tipos de IA en la industria",
             ["IA predictiva (mantenimiento)", "IA generativa (documentos, diseño)", "Visión por computadora (calidad)", "Cada tipo resuelve problemas distintos"],
             "En la industria hay tres grandes tipos de IA con aplicaciones específicas. La predictiva anticipa fallos, la generativa crea contenido y la visión por computadora automatiza el control de calidad."),
            ("Casos de uso probados",
             ["Mantenimiento predictivo de equipos", "Optimización energética de plantas", "Control de calidad automatizado", "Asistentes para documentación"],
             "Los casos de uso más maduros son el mantenimiento predictivo, la optimización energética y el control de calidad automatizado. Todos han demostrado retornos medibles en plantas industriales."),
            ("Datos: el combustible de la IA",
             ["Sensores e IoT en máquinas", "Datos históricos de producción", "Calidad y limpieza son críticas", "Sin datos buenos no hay IA útil"],
             "La IA depende completamente de la calidad de los datos disponibles. Sensores, históricos de producción y procesos limpios son el combustible; sin datos fiables ninguna IA ofrece resultados útiles."),
            ("Próximos pasos de este itinerario",
             ["Deep dive: IA en producción", "Deep dive: ética de la IA", "Cómo empezar en tu equipo"],
             "Esta introducción es el primer paso de un itinerario de tres lecciones. A continuación profundizaremos en la IA aplicada a producción y en las cuestiones éticas indispensables."),
        ],
    },
    "ia-en-produccion": {
        "title": "IA aplicada a producción",
        "lang": "es-ES",
        "tags": ["ai", "production", "es"],
        "slides": [
            ("IA aplicada a producción",
             ["Tres palancas de valor", "Cómo estructurar un proyecto", "Prerequisito: introducción a la IA"],
             "Esta lección aborda cómo pasar de la teoría a la práctica en producción. Cubriremos las tres palancas de valor principales y cómo estructurar un primer proyecto de IA industrial."),
            ("Mantenimiento predictivo",
             ["Sensores vibración y temperatura", "Modelos detectan anomalías", "Reducción paradas no planificadas", "Retorno típico: 20-30% ahorro"],
             "El mantenimiento predictivo combina sensores y modelos de machine learning para detectar anomalías antes del fallo. Las empresas reportan entre veinte y treinta por ciento de ahorro en paradas no planificadas."),
            ("Optimización energética",
             ["Modelos de consumo en tiempo real", "Ajuste automático de consignas", "Ahorro 5-15% en energía", "Reduce scope 2 directamente"],
             "La optimización energética basada en IA ajusta consignas en tiempo real según demanda y condiciones. Ahorros típicos de cinco a quince por ciento reducen directamente las emisiones de alcance dos."),
            ("Control de calidad con visión",
             ["Cámaras de alta resolución", "Modelos detectan defectos", "100% de piezas inspeccionadas", "Retroalimentación instantánea al operario"],
             "La visión por computadora permite inspeccionar el cien por cien de las piezas en tiempo real, algo imposible manualmente. La retroalimentación instantánea al operario mejora tanto calidad como productividad."),
            ("Cómo estructurar un proyecto",
             ["Empezar por un caso acotado", "Equipo mixto: negocio + datos", "Medir KPIs antes y después", "Escalar solo tras validar valor"],
             "Un proyecto de IA en producción debe empezar acotado, con un equipo mixto de negocio y datos. Medid KPIs antes y después, y escalad solo después de validar el valor real en planta."),
        ],
    },
    "etica-ia-empresa": {
        "title": "Ética de la IA en la empresa",
        "lang": "es-ES",
        "tags": ["ai", "ethics", "es"],
        "slides": [
            ("Ética de la IA en la empresa",
             ["Por qué la ética es clave", "Cinco principios aplicables", "Prerequisito: IA en producción"],
             "La última lección del itinerario aborda la ética de la IA, un tema transversal imprescindible. Una IA técnicamente buena pero mal gobernada puede generar más riesgos que beneficios."),
            ("Cinco principios fundamentales",
             ["Transparencia de los modelos", "Equidad y no discriminación", "Respeto de la privacidad", "Responsabilidad humana clara", "Robustez y seguridad"],
             "Cinco principios guían la IA responsable: transparencia, equidad, privacidad, responsabilidad humana y robustez técnica. Son aplicables tanto a modelos comprados como desarrollados internamente."),
            ("Sesgos: el riesgo invisible",
             ["Datos sesgados = modelos sesgados", "Ejemplo: decisiones de RRHH", "Auditoría regular necesaria", "Diversidad del equipo ayuda"],
             "Los sesgos son el riesgo más insidioso porque se propagan de los datos al modelo sin ser visibles. Auditorías regulares y equipos diversos son la mejor protección frente a decisiones discriminatorias."),
            ("IA y trabajo: el factor humano",
             ["La IA asiste, no sustituye", "Formación continua imprescindible", "Decisiones críticas siguen siendo humanas", "Comunicación transparente con empleados"],
             "La IA debe asistir a las personas, no sustituirlas en decisiones críticas. La formación continua y la comunicación transparente con los empleados son claves para una adopción sostenible."),
            ("Gobernanza de IA en la empresa",
             ["Comité de ética IA", "Inventario de modelos usados", "Evaluación de impacto por caso", "Formación obligatoria para usuarios"],
             "Una gobernanza eficaz combina un comité de ética, un inventario de modelos en uso, evaluaciones de impacto por caso y formación obligatoria para los usuarios. Es la base de una IA empresarial responsable."),
        ],
    },
}

ALL_GROUPS = {"A_safety_fr": GROUP_A_SAFETY_FR, "B_sustainability_en": GROUP_B_SUSTAINABILITY_EN, "C_ai_es": GROUP_C_AI_ES}


def _add_title_slide(prs: Presentation, title: str, bullets: list[str], notes: str) -> None:
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    if len(slide.placeholders) > 1:
        sub = slide.placeholders[1]
        sub.text = "\n".join(bullets)
    slide.notes_slide.notes_text_frame.text = notes


def _add_content_slide(prs: Presentation, title: str, bullets: list[str], notes: str) -> None:
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    body = slide.placeholders[1].text_frame
    body.text = bullets[0]
    for b in bullets[1:]:
        p = body.add_paragraph()
        p.text = b
        p.level = 0
    slide.notes_slide.notes_text_frame.text = notes


def build_deck(filename: str, spec: dict) -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slides = spec["slides"]
    # slide 0 = title, 1-3 = content, 4 = conclusion (all same layout except first)
    _add_title_slide(prs, *slides[0])
    for s in slides[1:]:
        _add_content_slide(prs, *s)
    out = OUT / f"{filename}.pptx"
    prs.save(out)
    return out


def main() -> None:
    index = []
    for group_key, decks in ALL_GROUPS.items():
        for deck_name, spec in decks.items():
            path = build_deck(deck_name, spec)
            index.append({
                "group": group_key,
                "file": path.name,
                "title": spec["title"],
                "language": spec["lang"],
                "tags": spec["tags"],
                "slide_count": len(spec["slides"]),
            })
            print(f"  ✓ {path.name}  ({spec['lang']}, {len(spec['slides'])} slides)")
    # Write manifest for automated tests
    import json as _json
    (OUT / "manifest.json").write_text(_json.dumps(index, indent=2, ensure_ascii=False), encoding="utf-8")
    print(f"\nGenerated {len(index)} decks in {OUT}")
    print(f"Manifest: {OUT / 'manifest.json'}")


if __name__ == "__main__":
    main()
